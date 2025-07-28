#!/usr/bin/env python3
import sys
import re
import os
import xml.etree.ElementTree as ET
import requests
from pptx import Presentation
# === Configuration ===
INVOKE_URL = "https://api.brev.dev/v1/chat/completions"
API_KEY    = "brev_api_-30Tr0kdRzKy3RTr4k9AtChl7TrL"
MODEL_NAME = "nvcf:nvidia/llama-3.1-nemotron-nano-8b-v1:dep-30TtAqxOULNBaELk6LSoLPKzNfV"
def run_inference(prompt: str, output_path: str = "GeneratedDeck.pptx") -> str:
    """
    Complete workflow: calls the model, saves XML, parses it, and generates PPTX.
    Returns the path to the generated PPTX file.
    """
    # Build the combined prompt with strong schema guardrail
    schema_instruction = (
        "Output ONLY valid XML in this exact schema (no comments, no prose, no extra tags!):\n"
        "<slide>\n"
        "  <title>…</title>\n"
        "  <bullet>…</bullet>\n"
        "  <bullet>…</bullet>\n"
        "  <bullet>…</bullet>\n"
        "  <bullet>…</bullet>\n"
        "</slide>"
    )
    final_prompt = f"{prompt}\n\n{schema_instruction}"
    
    # 1. Call the model
    slide_xml = _call_model(final_prompt)
    
    # 2. Save slide.xml
    save_xml(slide_xml, "slide.xml")
    print(":white_check_mark: slide.xml generated")
    
    # 3. Parse slide.xml
    slide_data = parse_slide_xml("slide.xml")
    
    # 4. Generate PPTX
    pptx_path = generate_pptx(slide_data, output_path)
    print(f":white_check_mark: PowerPoint saved to {pptx_path}")
    
    return pptx_path

def _call_model(prompt: str) -> str:
    """
    Calls the model and returns exactly the <slide>…</slide> XML.
    Raises if none is found.
    """
    headers = {
        "Authorization": f"Bearer {API_KEY}",
        "Content-Type":  "application/json",
        "Accept":        "application/json",
    }
    payload = {
        "model":      MODEL_NAME,
        "messages":   [{"role": "user", "content": prompt}],
        "max_tokens": 512,
        "temperature": 0.0,
        "top_p":       1.0,
        "stream":      False,
    }
    resp = requests.post(INVOKE_URL, json=payload, headers=headers)
    resp.raise_for_status()
    raw = resp.json()["choices"][0]["message"]["content"]
    # DEBUG: show raw to diagnose schema mismatches
    print("----- RAW MODEL OUTPUT -----")
    print(raw)
    print("----------------------------\n")
    # extract only the <slide>…</slide> portion
    match = re.search(r"<slide>.*?</slide>", raw, re.DOTALL)
    if not match:
        raise ValueError("No <slide>…</slide> found in model output")
    return match.group(0)
def save_xml(xml_str: str, path: str = "slide.xml"):
    with open(path, "w", encoding="utf-8") as f:
        f.write(xml_str)
def parse_slide_xml(xml_path: str):
    try:
        tree = ET.parse(xml_path)
        root = tree.getroot()
        title = root.find("title").text
        bullets = [b.text for b in root.findall("bullet")]
        return {"title": title, "bullets": bullets}
    except ET.ParseError as e:
        # If XML parsing fails, try to read and clean the file manually
        with open(xml_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Simple regex-based parsing as fallback
        import re
        title_match = re.search(r'<title>(.*?)</title>', content, re.DOTALL)
        bullet_matches = re.findall(r'<bullet>(.*?)</bullet>', content, re.DOTALL)
        
        if title_match and bullet_matches:
            title = title_match.group(1).strip()
            bullets = [bullet.strip() for bullet in bullet_matches]
            return {"title": title, "bullets": bullets}
        else:
            raise ValueError(f"Could not parse XML content: {e}")
def generate_pptx(slide_data: dict, output_path: str = "GeneratedDeck.pptx"):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = slide_data["title"]
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for txt in slide_data["bullets"]:
        p = tf.add_paragraph()
        p.text = txt
        p.level = 0
    prs.save(output_path)
    return output_path
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python generate_slide.py \"<your slide prompt>\"", file=sys.stderr)
        sys.exit(1)
    
    # Get the user prompt
    user_prompt = " ".join(sys.argv[1:])
    
    # Run the complete workflow
    pptx_path = run_inference(user_prompt)
    print(f"✅ Presentation generated successfully: {pptx_path}")
