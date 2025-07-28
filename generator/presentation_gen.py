from pptx import Presentation
import xml.etree.ElementTree as ET
import os

def xml_parser(xml_file_path):
    tree = ET.parse(xml_file_path)
    root = tree.getroot()

    slides_data = []
    
    for slide in root.findall('slide'):
        title = slide.find('title').text
        content = [point.text for point in slide.find('content').findall('point')]
        slides_data.append({'title': title, 'content': content})

    return slides_data

def presentation_gen(xml_file_path):
    # Parse the XML file
    slides_data = xml_parser(xml_file_path)
    
    prs = Presentation()
    for slide_info in slides_data:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = slide_info['title']
        content_shape = slide.placeholders[1]
        content_shape.text = ''

        for point in slide_info['content']:
            p = content_shape.text_frame.add_paragraph()
            p.text = point

    output_path = os.path.join(os.path.dirname(__file__), 'GeneratedDeck.pptx')
    prs.save(output_path)

    return output_path

if __name__ == '__main__':
    # For testing with the default slides.xml
    default_xml_path = os.path.join(os.path.dirname(__file__), 'slides.xml')
    presentation_gen(default_xml_path)